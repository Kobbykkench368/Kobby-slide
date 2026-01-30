import streamlit as st
from groq import Groq
from pptx import Presentation
from io import BytesIO

# 1. Setup Groq Client
# It will look for "GROQ_API_KEY" in your Streamlit Secrets
client = Groq(api_key=st.secrets["GROQ_API_KEY"])

def generate_slide_content(topic, num_slides):
    prompt = (
        f"Create a presentation outline for {num_slides} slides about '{topic}'. "
        "For each slide, follow this format exactly:\n"
        "TITLE: [Slide Title]\n"
        "CONTENT: [Point 1], [Point 2], [Point 3]\n"
    )
    
    completion = client.chat.completions.create(
        model="llama3-8b-8192",  # You can also use llama3-70b-8192
        messages=[
            {"role": "system", "content": "You are a professional slide generator."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.7,
    )
    return completion.choices[0].message.content

def create_pptx(topic, raw_text):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = topic
    slide.placeholders[1].text = "AI Generated with Groq & Python"

    # Split the AI response into slides
    slides_raw = raw_text.split("TITLE:")
    for slide_data in slides_raw[1:]:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        parts = slide_data.split("CONTENT:")
        slide.shapes.title.text = parts[0].strip()
        
        if len(parts) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            bullets = parts[1].strip().split(",")
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet.strip().replace('-', '').strip()

    # Save to memory buffer
    binary_output = BytesIO()
    prs.save(binary_output)
    return binary_output.getvalue()

# 2. Streamlit UI
st.set_page_config(page_title="Groq AI Slides", page_icon="⚡")
st.title("⚡ Groq AI Slide Generator")

with st.sidebar:
    st.info("Using Groq Cloud (Llama 3)")
    topic = st.text_input("Topic:", value="Digital Marketing in Ghana")
    count = st.slider("Slides:", 3, 10, 5)
    btn = st.button("Generate Now")

if btn:
    with st.spinner("Groq is working at lightning speed..."):
        try:
            content = generate_slide_content(topic, count)
            pptx_data = create_pptx(topic, content)
            
            st.success("Done!")
            st.download_button(
                label="📥 Download PowerPoint",
                data=pptx_data,
                file_name=f"{topic.replace(' ', '_')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            with st.expander("Show AI Raw Text"):
                st.write(content)
        except Exception as e:
            st.error(f"Something went wrong: {e}")
